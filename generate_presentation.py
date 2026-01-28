#!/usr/bin/env python3
"""
Generate a stylish project presentation (PPTX) with images.

Output: Iris_Recognition_Project_Presentation.pptx in the project root.

This script will:
- Create a modern title slide
- Add overview, features, architecture, datasets, model, demo, results, and roadmap slides
- Auto-embed best-available images from the repository (if present)

Requires: python-pptx, Pillow
"""

import os
import sys
from datetime import datetime

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TITLE = "Iris Recognition System - Secure Database Manager"
SUBTITLE = "Architecture, Security, and Operations"
OUTPUT = "Iris_Recognition_System_Presentation.pptx"

ACCENT = RGBColor(26, 115, 232)
DARK = RGBColor(33, 33, 33)
GRAY = RGBColor(97, 97, 97)


def ensure_deps():
    try:
        import pptx  # noqa: F401
        from PIL import Image  # noqa: F401
        return True
    except Exception:
        # Attempt on-the-fly install when run in venv
        try:
            import subprocess
            subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx', 'Pillow'])
            return True
        except Exception:
            return False


def first_existing(path_candidates):
    for p in path_candidates:
        if os.path.exists(p):
            return p
    return None


def add_full_width_image(slide, image_path, presentation_width, top_cm, height_cm=None):
    from pptx.util import Cm
    if not image_path:
        return None
    left = Cm(0.5)
    width = presentation_width - Cm(1.0)
    if height_cm is None:
        pic = slide.shapes.add_picture(image_path, left, Cm(top_cm), width=width)
    else:
        pic = slide.shapes.add_picture(image_path, left, Cm(top_cm), width=width, height=Cm(height_cm))
    return pic


def add_title_slide(prs: Presentation) -> None:
	slide = prs.slides.add_slide(prs.slide_layouts[0])
	slide.shapes.title.text = TITLE
	subtitle = slide.placeholders[1]
	subtitle.text = SUBTITLE
	subtitle.text_frame.paragraphs[0].font.color.rgb = GRAY


def add_bullets_slide(prs: Presentation, title: str, bullets: list) -> None:
	slide = prs.slides.add_slide(prs.slide_layouts[1])
	slide.shapes.title.text = title
	body = slide.shapes.placeholders[1].text_frame
	body.clear()
	for i, line in enumerate(bullets):
		p = body.add_paragraph() if i > 0 else body.paragraphs[0]
		p.text = line
		p.level = 0
		p.font.size = Pt(20)
		p.font.color.rgb = DARK


def add_two_column_slide(prs: Presentation, title: str, left: list, right: list) -> None:
	slide = prs.slides.add_slide(prs.slide_layouts[5])
	slide.shapes.title.text = title
	left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(4.2), Inches(4.5))
	right_box = slide.shapes.add_textbox(Inches(5.1), Inches(1.8), Inches(4.2), Inches(4.5))
	for idx, txt_list in enumerate((left, right)):
		tf = (left_box if idx == 0 else right_box).text_frame
		tf.clear()
		for i, line in enumerate(txt_list):
			p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
			p.text = line
			p.level = 0
			p.font.size = Pt(18)
			p.font.color.rgb = DARK


def build_presentation(output_path: str = OUTPUT) -> str:
	prs = Presentation()

	# Title
	add_title_slide(prs)

	# Overview
	add_bullets_slide(prs, "Project Overview", [
		"Biometric-driven access and voting using iris recognition",
		"SQLite-backed database manager with auditability",
		"Focus on security, integrity, and operational visibility",
	])

	# Architecture
	add_two_column_slide(prs, "High-Level Architecture", [
		"Iris capture + feature extraction (templates)",
		"Application layer: enrollment, auth, voting",
		"Data layer: persons, users, logs, voting, audit",
		"RBAC: admin/operator/viewer",
	], [
		"SQLite (WAL, FK enforcement)",
		"Hash-chained audit logs (tamper-evident)",
		"Access logs for forensics and monitoring",
		"Settings for configurable thresholds/policies",
	])

	# Database Schema
	add_bullets_slide(prs, "Key Tables", [
		"persons: identity, iris/face templates, metadata, activity",
		"users: login, role, TOTP, failed_attempts, lock_until",
		"access_logs: attempts with confidence, outcome, device/location",
		"voting_records: election_id, vote_hash, verification method",
		"audit_logs: event, actor, record_hash with prev_hash",
		"system_settings and model_versions",
	])

	# Security
	add_bullets_slide(prs, "Security Controls Implemented", [
		"Parameterized SQL queries (injection prevention)",
		"RBAC with least-privilege roles and soft-disable",
		"PBKDF2-SHA256 password hashing (salted, high iterations)",
		"Optional TOTP-based 2FA (RFC 6238)",
		"Account lockout policy",
		"Audit trail: SHA-256 hash chaining for tamper evidence",
		"Unique constraints and FK enforcement",
		"Basic input validation (email/phone)",
	])

	# Authentication & RBAC
	add_two_column_slide(prs, "Authentication & RBAC", [
		"Password hashing: PBKDF2-SHA256, constant-time verify",
		"Lockouts: failed_attempts + lock_until",
		"2FA: TOTP verification window handling",
	], [
		"RBAC roles: admin, operator, viewer",
		"Users optionally linked to persons",
		"Disable users via is_active flag",
	])

	# Audit and Logging
	add_bullets_slide(prs, "Auditability & Logging", [
		"Audit logs: record_hash = SHA-256(t, u, a, r, d, prev_hash)",
		"Access logs: result, score, device/location for forensics",
		"Indexes for efficient investigations",
	])

	# Workflows
	add_two_column_slide(prs, "Core Workflows", [
		"Enrollment: capture → serialize → store",
		"Authentication: match → RBAC check → log",
		"Voting: uniqueness check → vote_hash store",
	], [
		"Admin ops: user management, role assignment",
		"Maintenance: cleanup, backups, model versions",
		"Settings: configurable thresholds and policies",
	])

	# Roadmap
	add_bullets_slide(prs, "Roadmap", [
		"Field-level encryption for biometrics (Fernet/SQLCipher)",
		"Rate limiting and detection of anomalous attempts",
		"Comprehensive input validation and schema constraints",
		"API layer with JWT session security",
		"Automated backup signing and verification",
	])

	# Conclusion
	add_bullets_slide(prs, "Conclusion", [
		"Secure-by-design data layer for iris recognition",
		"Defense-in-depth: hashing, 2FA, lockouts, audits",
		"Clear paths for further hardening and scale",
	])

	prs.save(output_path)
	return os.path.abspath(output_path)


if __name__ == "__main__":
	path = build_presentation()
	print(f"Presentation generated: {path}")





